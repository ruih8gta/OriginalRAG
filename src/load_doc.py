import os
import sys
import configparser

import warnings
from langchain_community.document_loaders import TextLoader, UnstructuredWordDocumentLoader, PyPDFLoader, UnstructuredURLLoader, UnstructuredPowerPointLoader, UnstructuredExcelLoader
from langchain.text_splitter import CharacterTextSplitter

#ドキュメント読み込みライブラリ群
from docx import Document as word_doc  # Wordドキュメント用
from openpyxl import load_workbook  # Excel用
#from PyPDF2 import PdfReader  # PDF用
from pptx import Presentation  # PowerPoint用

from langchain_core.documents import Document
from pathlib import Path
import json
# Flask稼働時の設定ファイル読み込み
#config_path ="./config/config.ini"
# プログラム実行時の設定ファイル読み込み
config_path ="../config/config.ini"

config = configparser.ConfigParser()
config.read(config_path, encoding='utf-8') 

doc_path = config['doc']['DOC_DIR']
doc_admin_path = config['doc']['DOC_ADMIN_DIR']
doc_text_path = config['doc']['DOC_TEXT_DIR']
doc_db_path = config['doc']['DB_DIR']

current_dir = os.path.dirname(__file__)
parent_dir = os.path.dirname(current_dir)


#クラス
# Wordドキュメントからテキストを抽出するクラス
class Spec:
    def __init__(self,file_path):
        self.file_path = file_path

    def get_doc_text(self):
        document = word_doc(self.file_path)#選択したWordファイルのデータを取得
        spec=""
        spec+=self.get_text(document,spec)
        return spec

    def get_text(self,doc,gettext):
        for para in doc.paragraphs: #段落内のテキストを抽出
            gettext+=para.text+"\n"
        for table in doc.tables: #テーブル内のテキストを抽出
            for row in table.rows:
                for cell in row.cells:
                    for nested_paragraphs in cell.paragraphs: #入れ子になった段落がある場合
                        gettext+=nested_paragraphs.text+"\n"
                    for nested_tables in cell.tables:         #入れ子になったテーブルがある場合
                        gettext=self.get_text_from_table(nested_tables,gettext)
        return gettext

    def get_text_from_table(self,tbl,gettext):#入れ子になったテーブル内のテキストを再帰的に抽出
        for row in tbl.rows:
            for cell in row.cells:
                for nested_paragraphs in cell.paragraphs:
                    gettext+=nested_paragraphs.text+"\n"
                for nested_tables in cell.tables:
                    gettext=self.get_text_from_table(nested_tables,gettext)
        return gettext


# 指定したローカルフォルダに入っているファイルの読み込み
def Load_local_data(DIR_PATH,ADMIN_PATH,TEXT_PATH,DB_PATH):
    docs = []
    no=1
    for cur_dir, inc_dir, file_names in os.walk((DIR_PATH)):
        for file_name in file_names:
            try:    
                file_path = os.path.join(cur_dir, file_name)
                doc, content_type = Mimetype_Checker(file_path)
                if doc is not None:  # loaderがNoneでない場合のみ処理
                    file_info = {
                        "file_name": file_name,
                        "file_type": content_type,
                        "dir_path": cur_dir,
                        "file_path": file_path,
                    }

                    #doc =loader.load_and_split()
                    create_adminfile(no,doc,ADMIN_PATH,file_info,TEXT_PATH,DB_PATH) # docsをadmin_pathに保存する処理を実装
                    docs.extend(doc)
                    print("File path : " + file_path + " : File type :" + content_type)
                    no+=1
                else:
                    warnings.warn(f"Unsupported or unreadable file: {file_path}")
            except Exception as e:
                warnings.warn(f"Unreadable File: {file_path} - {str(e)}")
                pass
    return docs


# テキストローダー群
def extract_text_from_txt(filepath):
    """
    テキストファイルからコンテンツを抽出してLangChainのDocument形式に変換する
    
    引数:
        filepath: テキストファイルへのパス
        
    戻り値:
        ファイルの内容を含む単一のDocumentからなるリスト
    """
    # UTF-8エンコーディングでファイルを開き、エンコーディングエラーは無視する
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    # メタデータ付きのDocumentオブジェクトを作成して返す
    return [Document(page_content=content, metadata={"source": filepath, "type": "txt"})]

def extract_text_from_docx(filepath):
    """
    Wordドキュメントからコンテンツを抽出してLangChainのDocument形式に変換する
    
    引数:
        filepath: Wordドキュメントへのパス
        
    戻り値:
        ファイルの内容を含む単一のDocumentからなるリスト
    """
    # python-docxを使用してドキュメントを開く
    #doc = word_doc(filepath)
    # すべての段落テキストを改行で連結
    #content = '\n'.join([para.text for para in doc.paragraphs])
    word_instance = Spec(filepath)
    content = word_instance.get_doc_text()

    # メタデータ付きのDocumentオブジェクトを作成して返す
    return [Document(page_content=content, metadata={"source": filepath, "type": "docx"})]

def extract_text_from_xlsx(filepath):
    """
    Excelファイルからコンテンツを抽出してLangChainのDocument形式に変換する
    
    引数:
        filepath: Excelファイルへのパス
        
    戻り値:
        ファイルの内容を含む単一のDocumentからなるリスト
    """
    # 数式ではなく値を取得するためにdata_only=Trueでワークブックを読み込む
    wb = load_workbook(filepath, data_only=True)
    text = []
    # すべてのワークシートと行を反復処理
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            # すべてのセルを文字列に変換し、Noneの場合は空文字を使用
            line = [str(cell) if cell is not None else '' for cell in row]
            text.append('\t'.join(line))  # 同じ行のセルをタブで区切る
    # すべての行を改行で連結
    content = '\n'.join(text)
    # シート数を含むメタデータ付きのDocumentオブジェクトを作成して返す
    return [Document(page_content=content, metadata={"source": filepath, "type": "xlsx", "sheet_count": len(wb.worksheets)})]
"""
def extract_text_from_pdf(filepath):
    
    #PDFファイルからコンテンツを抽出してLangChainのDocument形式に変換する
    
    #引数:
    #    filepath: PDFファイルへのパス
        
    #戻り値:
    #    ファイルの内容を含む単一のDocumentからなるリスト
    
    # PyPDF2を使用してPDFを読み込む
    reader = PdfReader(filepath)
    text = []
    # 各ページからテキストを抽出
    for page in reader.pages:
        text.append(page.extract_text() or "")  # 抽出に失敗した場合は空文字を使用
    # すべてのページを改行で連結
    content = '\n'.join(text)
    # ページ数を含むメタデータ付きのDocumentオブジェクトを作成して返す
    return [Document(page_content=content, metadata={"source": filepath, "type": "pdf", "page_count": len(reader.pages)})]
"""
def extract_text_from_pptx(filepath):
    """
    PowerPointファイルからコンテンツを抽出してLangChainのDocument形式に変換する
    テキストを含むスライドごとに1つのDocumentオブジェクトを作成する
    
    引数:
        filepath: PowerPointファイルへのパス
        
    戻り値:
        テキストを含む各スライドに対応するDocumentオブジェクトのリスト
    """
    # プレゼンテーションを読み込む
    prs = Presentation(filepath)
    documents = []
    
    # 各スライドを個別に処理
    for i, slide in enumerate(prs.slides):
        slide_text = []
        # テキストプロパティを持つすべての図形からテキストを抽出
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        
        # スライドにテキスト内容がある場合のみドキュメントを作成
        if slide_text:
            # すべてのテキスト要素を改行で連結
            content = '\n'.join(slide_text)
            # スライド固有のメタデータを持つDocumentを作成
            doc = Document(
                page_content=content,
                metadata={
                    "source": filepath,
                    "type": "pptx",
                    "slide_number": i+1,  # 1から始まるスライド番号
                    "total_slides": len(prs.slides)
                }
            )
            documents.append(doc)
            
    return documents
# 読み込むファイル形式を読み取り、適したLoaderを使用
def Mimetype_Checker(file_path):
    PDF_MimeList = [".pdf"]
    TEXT_MimeList = [".txt"]
    HTML_MimeList = [".html"]
    Word_MimeList = [".doc", ".docx"]
    PPT_MimeList = [".ppt", ".pptx"]
    Excel_MimeList = [".xls", ".xlsx"]
    #ファイルの拡張子を取得
    content_type = os.path.splitext(file_path)[1].lower()

    if content_type in TEXT_MimeList:
        text_doc = extract_text_from_txt(file_path)
    elif content_type in PDF_MimeList:
        #text_doc = extract_text_from_pdf(file_path)
        text_doc = PyPDFLoader(file_path).load()
    elif content_type in HTML_MimeList:
        text_doc = UnstructuredURLLoader(file_path).load()
    elif content_type in Word_MimeList:
        text_doc = extract_text_from_docx(file_path)
    elif content_type in PPT_MimeList:
        text_doc = extract_text_from_pptx(file_path)
    elif content_type in Excel_MimeList:
        text_doc = extract_text_from_xlsx(file_path)
    else:
        warnings.warn(f"Unsupported filetype found. :" + file_path)
    return text_doc, content_type

def create_adminfile(no,doc,admin_path,file_info,text_path,db_path):
    #doc_text_file_path = os.path.join(text_path,file_info["file_path"]+".json")
    #admin_file_name = os.path.join(admin_path,file_info["file_path"]+".json")
    
    #file_info["file_path"]をフォルダ名ごとに分割し、_でつないで最後に.jsonをつける    
    #path_name = Path(file_info["file_path"])
    #parts = path_name.parts
    #joined = "_".join(part for part in parts if not part.endswith(":\\"))
    #print("joined : " + joined)
    
    
    #joined = "test.json"
    path_name = Path(file_info["file_name"])
    
    #doc_text_file_name = path_name.with_suffix(".json")
    doc_text_file_name = str(no) + ".json"
    admin_file_name = str(no) + ".json"
    doc_text_file_path = os.path.join(text_path,doc_text_file_name)
    admin_file_path = os.path.join(admin_path,admin_file_name)
    db_path = os.path.join(db_path, str(no))

    #print("admin_file_name : " + admin_file_path)
    #print("doc_text_file_path : " + doc_text_file_path)

    # docsをadmin_pathに保存する処理を実装
    doc_info={
        "id":no,
        "file_name": file_info["file_name"],
        "file_type": file_info["file_type"],
        "dir_path": file_info["dir_path"],
        "file_path": file_info["file_path"],
        "doc_text_file_path":doc_text_file_path,
        "vector_db_path":db_path
    }
    if not os.path.exists(db_path):
        os.makedirs(db_path, exist_ok=True)


    #doc_infoをjson形式 utf-8 で保存する処理を実装
    doc_info_json = json.dumps(doc_info,ensure_ascii=False, indent=4)
    with open(admin_file_path, "w", encoding="utf-8") as f:
        f.write(doc_info_json)
    
    # docをjson形式 utf-8 で保存する処理を実装
    
    # Documentを辞書に変換
    doc_dic = [{"page_content": d.page_content, "metadata": d.metadata} for d in doc]
    doc_json = json.dumps(doc_dic,ensure_ascii=False, indent=4)
    with open(doc_text_file_path, "w", encoding="utf-8") as f:
        f.write(doc_json)

if(__name__ == "__main__"):

    #filename="1.xlsx"
    #filename="2.pptx"
    #filename="3.pptx"
    #filename="4.pdf"
    #filename="5.docx"
    #filename="6.xlsx"


    """
    tmp_path =".\\tmp\\doc"

    doc, content_type = Mimetype_Checker(os.path.join(tmp_path,filename))

    with open(os.path.join(tmp_path,"test.txt"), "w", encoding="utf-8") as f:
        f.write(content_type)
        f.write("\n")
        f.write(str(doc))
    """
    docs = Load_local_data(os.path.join(parent_dir, doc_path),
                           os.path.join(parent_dir, doc_admin_path),
                           os.path.join(parent_dir, doc_text_path),
                           os.path.join(parent_dir,doc_db_path))
